"""Document generation service for the toolkit.

Provides two modes for creating DOCX and PPTX files, mirroring the
neo-back ``word_generator.py`` / ``presentation_filler.py`` patterns
but without Azure Blob Storage:

**Free-form mode** — build documents from scratch:
  - ``markdown_to_docx``  : Markdown → HTML → DOCX (via html2docx)
  - ``text_to_docx``      : Plain text → DOCX paragraphs (python-docx)
  - ``slides_to_pptx``    : List of slides → PPTX (python-pptx)

**Template-filling mode** — replace placeholders in existing templates:
  - ``fill_docx_template`` : Replace ``{{key}}`` in a .docx
  - ``fill_pptx_template`` : Replace ``<<key>>`` in a .pptx

All generated files are saved under ``tempfiles/`` and can be served
via the ``/files/{filename}/download`` route.
"""

import copy
import io
import logging
import re
import uuid
from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

TEMPFILES_DIR = Path(__file__).parent.parent / "tempfiles"
TEMPFILES_DIR.mkdir(exist_ok=True)

TEMPLATES_DIR = Path(__file__).parent.parent / "templates"

# Default DOCX templates — FR A4 branded formats
DOCX_TEMPLATE_OP = str(TEMPLATES_DIR / "FR_WORD TEMPLATE_OP_A4_2025.docx")
DOCX_TEMPLATE_WE_BLANK = str(
    TEMPLATES_DIR / "FR_WORD TEMPLATE_WE_A4_2025_BLANK PAGE.docx"
)

_SAFE_NAME_RE = re.compile(r"[^\w\s-]")


def _save_to_tempfiles(buffer: io.BytesIO, original_name: str) -> str:
    """Save a BytesIO buffer to tempfiles with a UUID prefix."""
    unique_name = f"{uuid.uuid4().hex}_{original_name}"
    dest = TEMPFILES_DIR / unique_name
    dest.write_bytes(buffer.getvalue())
    logger.info(f"Generated file saved: {dest}")
    return str(dest)


# ─── Free-form mode ──────────────────────────────────────────────────────────


def markdown_to_docx(
    markdown_text: str,
    title: str,
    subtitle: str = "",
) -> str:
    """Convert Markdown content to a DOCX file.

    Uses the ``markdown`` library to convert to HTML, then ``html2docx``
    to produce a DOCX.

    Returns:
        Absolute path of the generated DOCX in tempfiles/.
    """
    import markdown as md
    from html2docx import html2docx

    html = md.markdown(
        markdown_text,
        extensions=["tables", "fenced_code", "nl2br", "sane_lists"],
    )

    # html2docx returns a BytesIO containing a valid docx
    buf = html2docx(html, title=title)

    # Prepend subtitle if provided — reload and insert before first paragraph
    if subtitle:
        buf.seek(0)
        doc = DocxDocument(buf)
        p = doc.paragraphs[0] if doc.paragraphs else doc.add_paragraph()
        new_p = copy.deepcopy(p._element)
        new_p.clear()
        run = new_p.makeelement(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}r",
            {},
        )
        text_el = run.makeelement(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t",
            {},
        )
        text_el.text = subtitle
        run.append(text_el)
        new_p.append(run)
        p._element.addnext(new_p)

        buf = io.BytesIO()
        doc.save(buf)

    safe_title = _SAFE_NAME_RE.sub("", title)[:50].strip().replace(" ", "_")
    return _save_to_tempfiles(buf, f"{safe_title}.docx")


def text_to_docx(text: str, title: str) -> str:
    """Create a simple DOCX from plain text.

    Returns:
        Absolute path of the generated DOCX in tempfiles/.
    """
    doc = DocxDocument()
    doc.add_heading(title, level=1)

    for paragraph in text.split("\n"):
        stripped = paragraph.strip()
        if stripped:
            doc.add_paragraph(stripped)

    buf = io.BytesIO()
    doc.save(buf)

    safe_title = _SAFE_NAME_RE.sub("", title)[:50].strip().replace(" ", "_")
    return _save_to_tempfiles(buf, f"{safe_title}.docx")


def slides_to_pptx(
    slides: list[dict[str, str]],
    title: str,
) -> str:
    """Create a PPTX presentation from a list of slide data.

    Each slide dict should contain ``title`` and ``content`` keys.

    Returns:
        Absolute path of the generated PPTX in tempfiles/.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = f"{len(slides)} slides"

    # Content slides
    content_layout = prs.slide_layouts[1]
    for slide_data in slides:
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = slide_data.get("title", "")
        body = s.placeholders[1]
        tf = body.text_frame
        tf.clear()
        content = slide_data.get("content", "")
        for i, line in enumerate(content.split("\n")):
            if i == 0:
                tf.paragraphs[0].text = line
                tf.paragraphs[0].font.size = Pt(16)
            else:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)

    buf = io.BytesIO()
    prs.save(buf)

    safe_title = _SAFE_NAME_RE.sub("", title)[:50].strip().replace(" ", "_")
    return _save_to_tempfiles(buf, f"{safe_title}.pptx")


# ─── Template-filling mode ───────────────────────────────────────────────────


def _replace_placeholders_in_paragraph(
    paragraph,
    placeholders: dict[str, str],
    left_delim: str,
    right_delim: str,
) -> None:
    """Replace delimited placeholders in a paragraph, preserving runs."""
    full_text = paragraph.text
    new_text = full_text
    for key, value in placeholders.items():
        token = f"{left_delim}{key}{right_delim}"
        new_text = new_text.replace(token, value)

    if new_text == full_text:
        return

    for run in paragraph.runs:
        run.text = ""
    if paragraph.runs:
        paragraph.runs[0].text = new_text
    else:
        paragraph.add_run(new_text)


def fill_docx_template(
    placeholders: dict[str, str],
    output_name: str,
    template_path: str | None = None,
) -> str:
    """Fill a DOCX template by replacing ``{{key}}`` placeholders.

    Returns:
        Absolute path of the generated DOCX in tempfiles/.
    """
    resolved = template_path or DOCX_TEMPLATE_OP
    if not Path(resolved).exists():
        raise FileNotFoundError(f"Template not found: {resolved}")

    doc = DocxDocument(resolved)

    for paragraph in doc.paragraphs:
        _replace_placeholders_in_paragraph(paragraph, placeholders, "{{", "}}")

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    _replace_placeholders_in_paragraph(
                        paragraph, placeholders, "{{", "}}"
                    )

    buf = io.BytesIO()
    doc.save(buf)

    safe_name = _SAFE_NAME_RE.sub("", output_name)[:50].strip().replace(" ", "_")
    return _save_to_tempfiles(buf, f"{safe_name}.docx")


def fill_pptx_template(
    template_path: str,
    placeholders: dict[str, str],
    output_name: str,
) -> str:
    """Fill a PPTX template by replacing ``<<key>>`` placeholders.

    Returns:
        Absolute path of the generated PPTX in tempfiles/.
    """
    if not Path(template_path).exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                _replace_placeholders_in_paragraph(
                    paragraph, placeholders, "<<", ">>"
                )

    buf = io.BytesIO()
    prs.save(buf)

    safe_name = _SAFE_NAME_RE.sub("", output_name)[:50].strip().replace(" ", "_")
    return _save_to_tempfiles(buf, f"{safe_name}.pptx")
