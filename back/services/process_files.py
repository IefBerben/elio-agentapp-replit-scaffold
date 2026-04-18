"""File processing service for the toolkit.

Provides local file extraction and conditional summarization, mirroring the
neo-back `print_or_summarize` function without Azure Blob Storage dependencies.

Files are expected to exist on the local filesystem (e.g. under tempfiles/).

Supported formats:
  - .pdf     (PyMuPDF primary, Azure Document Intelligence fallback for scanned PDFs)
  - .docx    (python-docx)
  - .pptx    (python-pptx)
  - .xlsx    (pandas)
  - .csv     (pandas)
  - .txt, .md, .html, .xml, .json, .svg, .srt, .py, .js  (plain read)
  - .mp3, .m4a, .wav, .ogg, .webm  (Whisper transcription)

Azure Document Intelligence (optional):
  Set AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT in .env to enable OCR-grade PDF extraction.
  Used automatically as fallback when PyMuPDF extracts fewer than 4 000 characters
  (typically scanned or image-only PDFs). Set FORCE_DOC_INTELLIGENCE=true to always
  use it regardless of PyMuPDF output.

  RBAC required: "Cognitive Services User" role on the Document Intelligence resource.
"""

import asyncio
import json
import logging
import os
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Literal

import fitz  # PyMuPDF
import pandas as pd
from azure.identity import DefaultAzureCredential, get_bearer_token_provider
from docx import Document as DocxDocument
from langchain_community.document_loaders import AzureAIDocumentIntelligenceLoader
from langchain_core.messages import SystemMessage
from openai import AzureOpenAI
from pptx import Presentation

from services.llm_config import get_llm

logger = logging.getLogger(__name__)

_executor = ThreadPoolExecutor(max_workers=4)

DEFAULT_THRESHOLD = 10_000
"""Characters above which content is summarized instead of returned raw."""

_MIN_PDF_CONTENT_LENGTH = 4_000
"""Below this char count, PyMuPDF output is considered poor and DI fallback is tried."""

_TEXT_EXTENSIONS = {
    ".txt", ".md", ".html", ".xml", ".json", ".svg", ".srt",
    ".py", ".js", ".ts", ".css",
}

_AUDIO_EXTENSIONS = {".mp3", ".m4a", ".wav", ".ogg", ".webm"}
"""Audio formats supported via Whisper transcription."""


# ─── Extraction helpers (sync, run in thread pool) ───────────────────────────


def _extract_pdf_pymupdf(file_path: str) -> str:
    """Extract text from a PDF using PyMuPDF (fast, text-based PDFs)."""
    doc = fitz.open(file_path)
    parts = []
    for i, page in enumerate(doc, start=1):
        text = page.get_text()
        if text.strip():
            parts.append(f"## Page {i}\n{text}")
    doc.close()
    return "\n".join(parts)


def _extract_pdf_doc_intelligence(file_path: str) -> str:
    """Extract text from a PDF using Azure Document Intelligence (OCR-grade).

    Uses ``prebuilt-layout`` model in ``page`` mode.
    Authenticates via DefaultAzureCredential.

    Requires AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT to be set in .env.
    RBAC: "Cognitive Services User" role on the Document Intelligence resource.
    """
    endpoint = os.getenv("AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT")
    if not endpoint:
        raise RuntimeError(
            "AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT is not set. "
            "Add it to .env to enable Document Intelligence OCR."
        )

    loader = AzureAIDocumentIntelligenceLoader(
        api_endpoint=endpoint,
        azure_credential=DefaultAzureCredential(),
        api_model="prebuilt-layout",
        mode="page",
        file_path=file_path,
    )
    docs = loader.load()
    parts = []
    for i, doc in enumerate(docs, start=1):
        page_num = doc.metadata.get("page", i)
        if doc.page_content.strip():
            parts.append(f"## Page {page_num}\n{doc.page_content}")
    return "\n".join(parts)


def _extract_pdf(file_path: str) -> str:
    """Extract text from a PDF, with Document Intelligence fallback for scanned PDFs.

    Strategy:
    - If FORCE_DOC_INTELLIGENCE=true → use Document Intelligence directly.
    - Otherwise, try PyMuPDF first. If the result is too short
      (< _MIN_PDF_CONTENT_LENGTH chars) AND DI endpoint is configured → fall back.
    """
    if os.getenv("FORCE_DOC_INTELLIGENCE") == "true":
        logger.info(f"Using Document Intelligence for {file_path} (forced)")
        return _extract_pdf_doc_intelligence(file_path)

    content = _extract_pdf_pymupdf(file_path)

    if (
        len(content) < _MIN_PDF_CONTENT_LENGTH
        and os.getenv("AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT")
    ):
        logger.info(
            f"PyMuPDF extracted only {len(content)} chars from {file_path}, "
            "falling back to Document Intelligence OCR"
        )
        try:
            di_content = _extract_pdf_doc_intelligence(file_path)
            if di_content:
                return di_content
        except Exception as exc:
            logger.warning(f"Document Intelligence fallback failed: {exc}")

    return content


def _extract_docx(file_path: str) -> str:
    """Extract text from a DOCX file using python-docx."""
    doc = DocxDocument(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _slide_text(slide) -> list[str]:
    """Extract non-empty text lines from a PPTX slide."""
    lines = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            lines.extend(
                para.text.strip()
                for para in shape.text_frame.paragraphs
                if para.text.strip()
            )
    return lines


def _extract_pptx(file_path: str) -> str:
    """Extract text from a PPTX file using python-pptx."""
    prs = Presentation(file_path)
    parts = [
        f"## Slide {i}\n" + "\n".join(lines)
        for i, slide in enumerate(prs.slides, start=1)
        if (lines := _slide_text(slide))
    ]
    return "\n".join(parts)


def _extract_excel(file_path: str) -> str:
    """Extract content from an XLSX file using pandas."""
    xls = pd.ExcelFile(file_path)
    parts = []
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        parts.append(f"## Sheet: {sheet_name}\n{df.to_string(index=False, max_rows=200)}")
    return "\n\n".join(parts)


def _extract_csv(file_path: str) -> str:
    """Extract content from a CSV file using pandas."""
    df = pd.read_csv(file_path, nrows=500)
    return df.to_string(index=False)


def _extract_text(file_path: str) -> str:
    """Read a plain text file (UTF-8 with latin-1 fallback)."""
    try:
        with open(file_path, encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(file_path, encoding="latin-1") as f:
            return f.read()


def _extract_audio(file_path: str) -> str:
    """Transcribe an audio file using Azure OpenAI Whisper deployment.

    Uses the same Azure OpenAI endpoint as LLM calls. The Whisper deployment
    name defaults to ``"whisper"`` and can be overridden via the
    ``WHISPER_DEPLOYMENT_NAME`` environment variable.
    """
    config_path = Path(__file__).parent / "config_llms.json"
    with open(config_path, encoding="utf-8") as f:
        cfg = json.load(f)

    endpoint = cfg["azure_openai"]["endpoint"]
    api_version = cfg["azure_openai"]["api_version"]
    deployment = os.getenv(
        "WHISPER_DEPLOYMENT_NAME",
        cfg.get("defaults", {}).get("whisper", "whisper"),
    )

    token_provider = get_bearer_token_provider(
        DefaultAzureCredential(),
        "https://cognitiveservices.azure.com/.default",
    )

    client = AzureOpenAI(
        azure_endpoint=endpoint,
        api_version=api_version,
        azure_ad_token_provider=token_provider,
    )

    logger.info(f"Transcribing audio file: {file_path} (deployment={deployment})")

    with open(file_path, "rb") as audio_file:
        result = client.audio.transcriptions.create(
            model=deployment,
            file=audio_file,
        )

    transcript = result.text if hasattr(result, "text") else str(result)
    logger.info(
        f"Transcription complete: {len(transcript)} chars from {file_path}"
    )
    return transcript


def _extract_sync(file_path: str) -> str:
    """Dispatch extraction by file extension (blocking)."""
    _, ext = os.path.splitext(file_path.lower())
    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext == ".docx":
        return _extract_docx(file_path)
    if ext == ".pptx":
        return _extract_pptx(file_path)
    if ext == ".xlsx":
        return _extract_excel(file_path)
    if ext == ".csv":
        return _extract_csv(file_path)
    if ext in _TEXT_EXTENSIONS:
        return _extract_text(file_path)
    if ext in _AUDIO_EXTENSIONS:
        return _extract_audio(file_path)
    raise ValueError(f"Unsupported file format: {ext}")


async def extract_file_content(file_path: str) -> str:
    """Asynchronously extract text from a local file.

    Runs the blocking extraction in a thread pool to keep the event loop free.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(_executor, _extract_sync, file_path)


# ─── Summarization ───────────────────────────────────────────────────────────

_SUMMARIZE_PROMPT = """You are a concise document summarizer.

Summarize the following document content in {language}, preserving all key \
information, names, dates, and facts. Be thorough but concise.
Output only the summary, no preamble.

Document:
{content}"""


async def summarize_text(
    text: str,
    language: Literal["FR", "EN"] = "FR",
    model: str = "",
) -> str:
    """Summarize a text using the configured LLM."""
    llm = get_llm(model)
    lang_name = "French" if language == "FR" else "English"
    prompt = _SUMMARIZE_PROMPT.format(content=text[:100_000], language=lang_name)
    response = await llm.ainvoke([SystemMessage(content=prompt)])
    return str(response.content).strip()


# ─── Main interface ───────────────────────────────────────────────────────────


async def print_or_summarize(
    file_path: str,
    threshold: int = DEFAULT_THRESHOLD,
    language: Literal["FR", "EN"] = "FR",
    model: str = "",
) -> str:
    """Extract content from a local file, summarizing if it exceeds the threshold.

    Args:
        file_path: Local path to the file.
        threshold: Character count above which content is summarized.
        language: Summary language if summarization is triggered ("FR" or "EN").
        model: LLM model name for summarization (empty string = default model).

    Returns:
        File content prefixed with its filename, or a summary if too long.
    """
    filename = os.path.basename(file_path)
    content = await extract_file_content(file_path)

    if len(content) <= threshold:
        return f"=== {filename} ===\n{content}"

    logger.info(f"Content too long ({len(content)} chars), summarizing {filename}…")
    summary = await summarize_text(content, language=language, model=model)
    return f"=== {filename} (summary) ===\n{summary}"
